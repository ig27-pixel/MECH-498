"""Export RoboRoll pitch deck to PowerPoint using python-pptx.

    pip install python-pptx pillow
Run from the repository root:
    python Project/pitch_deck/export_pitch_deck.py
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ─────────────────────────────────────────────────────────────────────
ROOT         = Path(__file__).resolve().parent
IMAGE_DIR    = ROOT / "images"
DYNAMICS_DIR = ROOT.parent / "section_4_dynamics"
OUT          = ROOT / "roboroll_pitch_deck.pptx"

# ── slide geometry (16 : 9 widescreen) ───────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)
PX = Inches(0.62)        # horizontal padding
PY = Inches(0.50)        # vertical padding
LW = Inches(6.1)         # left-column width
RX = Inches(7.15)        # right-column image left edge
RW = Inches(5.6)         # right-column width
RY = Inches(0.55)        # right-column top
RH = Inches(6.4)         # right-column height

# ── brand palette ─────────────────────────────────────────────────────────────
BLUE   = RGBColor(0x1F, 0x6F, 0xEB)
GREEN  = RGBColor(0x1B, 0x8F, 0x5A)
ORANGE = RGBColor(0xD6, 0x5F, 0x21)
DARK   = RGBColor(0x11, 0x18, 0x27)
INK    = RGBColor(0x17, 0x20, 0x2A)
MUTED  = RGBColor(0x5B, 0x66, 0x73)
PAPER  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HERO_TINT = RGBColor(0x8F, 0xC0, 0xFF)
HERO_SUB  = RGBColor(0xE5, 0xED, 0xF7)

FONT = "Arial"


# ── slide data ────────────────────────────────────────────────────────────────
# bullets  : list[str]                  → bullet-list left column
# tiles    : list[(title, body)]        → stacked tiles, green accent
# steps    : list[(title, body)]        → stacked steps, blue accent (numbered)
# risks    : list[(title, body)]        → 2×2 grid, orange accent, no image
# image    : str                        → single right-column image
# images   : list[str]                  → two stacked right-column images
# image_dir: "dynamics"                 → resolve from DYNAMICS_DIR instead
# kind     : "hero" | "dark"            → special full-slide layouts

SLIDES: list[dict] = [

    # ── title ─────────────────────────────────────────────────────────────────
    {
        "kind": "hero",
        "title": "RoboRoll Coatings",
        "subtitle": "Robotic wall painting for faster, safer, more repeatable interior finishing.",
        "image": "mockup_robot_hero.png",
    },

    # ── live pitch ────────────────────────────────────────────────────────────
    {
        "eyebrow": "1. The Problem",
        "title": "Interior painting is still a manual bottleneck.",
        "bullets": [
            "Repetitive wall coverage consumes crew hours.",
            "Quality changes with fatigue, reach, lighting, and setup consistency.",
            "Labor bottlenecks delay final construction turnover.",
        ],
        "image": "mockup_problem_scene.png",
    },
    {
        "eyebrow": "2. Why It Matters",
        "title": "Construction needs focused automation that helps crews now.",
        "tiles": [
            ("Common task",    "Painting shows up across apartments, commercial spaces, and new builds."),
            ("Repeatable work","Large wall areas create predictable paths that a robot can execute."),
            ("Schedule value", "Fewer rework passes mean more predictable finishing timelines."),
        ],
        "image": "mockup_market_segments.png",
    },
    {
        "eyebrow": "3. The Solution",
        "title": "A portable 4-DOF robot that follows planned wall paths.",
        "bullets": [
            "Base yaw aims the arm around the room.",
            "Shoulder, elbow, and wrist place the nozzle on target.",
            "Inverse kinematics converts wall points into smooth joint motion.",
        ],
        "image": "mockup_robot_architecture.png",
    },
    {
        "eyebrow": "4. Product Demo",
        "title": "The simulated robot already paints across two walls.",
        "bullets": [
            "Wall 1: smiley face — outline, eyes, and smile arc.",
            "Wall 2: five horizontal color stripes.",
            "Curved paths prove more than straight-line coverage.",
        ],
        "images": ["mockup_product_smiley.png", "mockup_product_lines.png"],
    },
    {
        "eyebrow": "5. Why We Win",
        "title": "RoboRoll starts with one focused construction task.",
        "tiles": [
            ("Focused wedge",   "Interior wall finishing is narrow enough for a first product."),
            ("Technical proof", "FK, IK, trajectory generation, and dynamics are already simulated."),
            ("The ask",         "Seed funding turns the validated simulation into a portable prototype."),
        ],
        "image": "mockup_roadmap.png",
    },

    # ── transitions ───────────────────────────────────────────────────────────
    {
        "kind": "dark",
        "title": "Thank You",
        "subtitle": "RoboRoll Coatings — robotic painting for the next generation of construction crews.",
    },
    {
        "kind": "dark",
        "eyebrow": "Additional Materials",
        "title": "Leave-Behind Slides",
        "subtitle": "The following slides provide the fuller technical, market, and roadmap story.",
    },

    # ── leave-behind ──────────────────────────────────────────────────────────
    {
        "eyebrow": "A1. Beachhead Market",
        "title": "Start with multifamily apartment turns.",
        "tiles": [
            ("Why first",    "Similar room layouts create repeatable painting paths."),
            ("Buyer pain",   "Turnover schedules are tight; labor availability affects move-in dates."),
            ("Next markets", "Commercial interiors and new-build finishing once setup is proven."),
        ],
        "image": "mockup_market_segments.png",
    },
    {
        "eyebrow": "A2. Customer Workflow",
        "title": "RoboRoll helps one operator supervise repeatable wall coverage.",
        "steps": [
            ("1  Position",  "Roll the robot into the room near the target wall."),
            ("2  Calibrate", "Register the wall plane and reachable painting area."),
            ("3  Paint",     "Execute planned coverage paths with human supervision."),
            ("4  Verify",    "Inspect coverage, refill, and move to the next room."),
        ],
        "image": "mockup_robot_workflow.png",
    },
    {
        "eyebrow": "A3. Robot Architecture",
        "title": "The first design is intentionally simple.",
        "bullets": [
            "4 revolute joints: base yaw, shoulder pitch, elbow pitch, wrist pitch.",
            "500 mm base height with 700 mm and 500 mm arm links.",
            "200 mm nozzle offset for reaching the wall surface.",
        ],
        "image": "mockup_nozzle_closeup.png",
    },
    {
        "eyebrow": "A4. Kinematics Proof",
        "title": "The robot can turn wall targets into joint commands.",
        "bullets": [
            "Forward kinematics verifies the end-effector pose.",
            "Inverse kinematics solves reachable wall targets.",
            "Joint limits prevent unrealistic poses during planning.",
        ],
        "image": "mockup_calibration.png",
    },
    {
        "eyebrow": "A5. Demo Path Details",
        "title": "The demo combines coverage and detail work.",
        "bullets": [
            "Hundreds of rendered frames across both walls.",
            "Smoothstep interpolation between IK waypoints.",
            "Lift motions prevent unwanted diagonal paint marks.",
        ],
        "image": "mockup_path_planning.png",
    },
    {
        "eyebrow": "A6. Dynamics Proof",
        "title": "The simulation includes physical reasoning.",
        "bullets": [
            "Passive motion checks kinetic, potential, and total energy behavior.",
            "PD plus gravity feedforward controls motion to a painting target.",
            "Torque and velocity plots support physically reasonable motion.",
        ],
        "images": ["dynamics_passive.png", "dynamics_controlled.png"],
        "image_dir": "dynamics",
    },
    {
        "eyebrow": "A7. Development Roadmap",
        "title": "The next milestones move from simulation to hardware.",
        "steps": [
            ("End effector", "Prototype paint delivery and nozzle control."),
            ("Calibration",  "Add wall sensing and room registration."),
            ("Testing",      "Measure repeatability on drywall panels."),
            ("Workflow",     "Build the operator setup process."),
        ],
        "image": "mockup_roadmap.png",
    },
    {
        "eyebrow": "A8. Risks and Mitigations",
        "title": "The key risks are practical and testable.",
        "risks": [
            ("Surface quality",   "Test nozzle spacing, flow rate, and overlap on drywall panels."),
            ("Wall localization", "Add sensing and calibration before each room run."),
            ("Setup time",        "Design a simple operator workflow with repeatable presets."),
            ("Safety",            "Limit speed, add supervised operation, and define a work envelope."),
        ],
    },
    {
        "eyebrow": "A9. Funding Use",
        "title": "Funding turns the validated simulation into a prototype.",
        "steps": [
            ("Arm",        "Build a first portable mechanism."),
            ("Paint tool", "Integrate nozzle, flow, and cleanup hardware."),
            ("Sensing",    "Add wall plane detection and calibration."),
            ("Testing",    "Run repeatability and finish-quality trials."),
        ],
    },
    {
        "kind": "dark",
        "title": "A focused wedge into construction automation.",
        "subtitle": (
            "RoboRoll begins with repeatable interior painting, proves value in "
            "controlled room-scale tasks, and expands toward broader finishing automation."
        ),
    },
]


# ── low-level helpers ─────────────────────────────────────────────────────────

def _img_path(name: str, data: dict) -> Path:
    if data.get("image_dir") == "dynamics":
        return DYNAMICS_DIR / name
    return IMAGE_DIR / name


def _img_dims(path: Path) -> tuple[int, int]:
    with Image.open(path) as im:
        return im.size


def _add_picture_fitted(slide, path: Path, x, y, max_w, max_h) -> None:
    """Add an image aspect-ratio-fitted and centred within a bounding box."""
    iw, ih = _img_dims(path)
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path),
                             x + (max_w - w) // 2,
                             y + (max_h - h) // 2,
                             w, h)


def _add_rect(slide, x, y, w, h, fill: RGBColor,
              transparency: float = 0.0, no_line: bool = True):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if transparency:
        shape.fill.fore_color.transparency = transparency
    if no_line:
        shape.line.fill.background()
    return shape


def _send_to_back(slide, shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _txb(slide, text: str, x, y, w, h, *,
         size: float, color: RGBColor,
         bold: bool = False,
         italic: bool = False,
         align: PP_ALIGN = PP_ALIGN.LEFT,
         caps: bool = False,
         space_before: float = 0):
    """Add a single-paragraph text box."""
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text.upper() if caps else text
    f = run.font
    f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
    f.color.rgb = color
    return box


# ── left-column header ────────────────────────────────────────────────────────

def _header(slide, eyebrow: str, title: str):
    """Draw eyebrow + title; return the Y coordinate where content starts."""
    _txb(slide, eyebrow, PX, PY, LW, Inches(0.32),
         size=11, color=BLUE, bold=True, caps=True)
    _txb(slide, title, PX, PY + Inches(0.35), LW, Inches(1.45),
         size=30, color=INK, bold=True)
    return PY + Inches(1.88)


# ── right-column image ────────────────────────────────────────────────────────

def _right_image(slide, data: dict) -> None:
    if "images" in data:
        per_h = (RH - Inches(0.16)) // 2
        for i, name in enumerate(data["images"]):
            _add_picture_fitted(slide, _img_path(name, data),
                                RX, RY + i * (per_h + Inches(0.16)), RW, per_h)
    elif "image" in data:
        _add_picture_fitted(slide, _img_path(data["image"], data),
                            RX, RY, RW, RH)


# ── tile drawing ──────────────────────────────────────────────────────────────

def _tile(slide, label: str, body: str, x, y, w, h, accent: RGBColor) -> None:
    """One tile: coloured top bar, bold label, muted body."""
    _add_rect(slide, x, y, w, Inches(0.055), accent)
    box = slide.shapes.add_textbox(x, y + Inches(0.08), w, h - Inches(0.09))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = label
    r1.font.name, r1.font.size = FONT, Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = INK

    p2 = tf.add_paragraph()
    p2.space_before = Pt(5)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name, r2.font.size = FONT, Pt(13)
    r2.font.color.rgb = MUTED


def _stacked_tiles(slide, items: list[tuple[str, str]],
                   content_y, accent: RGBColor) -> None:
    """Stack tiles vertically in the left column."""
    n = len(items)
    gap = Inches(0.13)
    avail = SH - Inches(0.28) - content_y
    tile_h = (avail - gap * (n - 1)) / n
    for i, (label, body) in enumerate(items):
        _tile(slide, label, body,
              PX, content_y + i * (tile_h + gap), LW, tile_h, accent)


def _grid2x2(slide, items: list[tuple[str, str]],
             content_y, accent: RGBColor) -> None:
    """2×2 tile grid spanning the full content width (no image column)."""
    full_w = SW - 2 * PX
    col_w  = (full_w - Inches(0.18)) / 2
    avail  = SH - Inches(0.28) - content_y
    row_h  = (avail - Inches(0.14)) / 2
    for i, (label, body) in enumerate(items):
        col, row = i % 2, i // 2
        _tile(slide, label, body,
              PX + col * (col_w + Inches(0.18)),
              content_y + row * (row_h + Inches(0.14)),
              col_w, row_h, accent)


# ── slide builders ────────────────────────────────────────────────────────────

def _hero(slide, data: dict) -> None:
    pic = slide.shapes.add_picture(
        str(IMAGE_DIR / data["image"]), 0, 0, SW, SH)
    _send_to_back(slide, pic)

    overlay = _add_rect(slide, 0, 0, SW, SH,
                        RGBColor(0x0A, 0x12, 0x20), transparency=0.28)
    _send_to_back(slide, overlay)

    ty = Inches(4.15)
    _txb(slide, "Startup Pitch", PX, ty, Inches(8), Inches(0.34),
         size=11, color=HERO_TINT, bold=True, caps=True)
    _txb(slide, data["title"], PX, ty + Inches(0.38), Inches(8.5), Inches(1.15),
         size=52, color=WHITE, bold=True)
    _txb(slide, data["subtitle"], PX, ty + Inches(1.6), Inches(7.8), Inches(1.0),
         size=22, color=HERO_SUB)


def _dark(slide, data: dict) -> None:
    _set_bg(slide, DARK)
    cy = Inches(2.3) if not data.get("eyebrow") else Inches(1.9)
    if data.get("eyebrow"):
        _txb(slide, data["eyebrow"],
             PX, cy - Inches(0.42), SW - 2 * PX, Inches(0.34),
             size=11, color=HERO_TINT, bold=True, caps=True,
             align=PP_ALIGN.CENTER)
    _txb(slide, data["title"],
         PX, cy, SW - 2 * PX, Inches(1.2),
         size=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        _txb(slide, data["subtitle"],
             PX, cy + Inches(1.3), SW - 2 * PX, Inches(1.8),
             size=22, color=HERO_SUB, align=PP_ALIGN.CENTER)


def _bullets(slide, data: dict) -> None:
    _set_bg(slide, PAPER)
    y = _header(slide, data["eyebrow"], data["title"])

    box = slide.shapes.add_textbox(PX + Inches(0.05), y,
                                   LW - Inches(0.05), SH - y - Inches(0.28))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True

    for i, text in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(14 if i > 0 else 0)
        run = p.add_run()
        run.text = "•  " + text   # bullet + en-space × 2
        run.font.name  = FONT
        run.font.size  = Pt(20)
        run.font.color.rgb = MUTED

    _right_image(slide, data)


def _tiles(slide, data: dict, key: str, accent: RGBColor) -> None:
    _set_bg(slide, PAPER)
    y = _header(slide, data["eyebrow"], data["title"])
    _stacked_tiles(slide, data[key], y, accent)
    _right_image(slide, data)


def _risks(slide, data: dict) -> None:
    _set_bg(slide, PAPER)
    y = _header(slide, data["eyebrow"], data["title"])
    _grid2x2(slide, data["risks"], y, ORANGE)


# ── presentation assembly ─────────────────────────────────────────────────────

def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[-1]


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    layout = _blank_layout(prs)

    for data in SLIDES:
        slide = prs.slides.add_slide(layout)
        kind  = data.get("kind")

        if kind == "hero":
            _hero(slide, data)
        elif kind == "dark":
            _dark(slide, data)
        elif "tiles" in data:
            _tiles(slide, data, "tiles", GREEN)
        elif "steps" in data:
            _tiles(slide, data, "steps", BLUE)
        elif "risks" in data:
            _risks(slide, data)
        else:
            _bullets(slide, data)

    if OUT.exists():
        OUT.unlink()
    prs.save(str(OUT))
    print(f"Saved {len(SLIDES)} slides: {OUT}")


if __name__ == "__main__":
    build_pptx()
